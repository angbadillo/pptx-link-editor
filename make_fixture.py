"""Genera un .pptx de prueba con enlaces de datos EXTERNOS (gráfico -> xlsx)."""
import zipfile, shutil, re, sys
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

base = "fixture_base.pptx"
out = "fixture_external.pptx"

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])

data = CategoryChartData()
data.categories = ["Q1", "Q2", "Q3"]
data.add_series("Ventas", (10.0, 20.0, 30.0))
slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                       Inches(1), Inches(1.5), Inches(6), Inches(4), data)
prs.save(base)

# Reescribir el .rels del gráfico para que el libro de datos sea EXTERNO,
# imitando un .pptx cuyo xlsx vive en otra ruta.
ext_target = "file:///C:\\Users\\old\\Reports\\Q1\\datos_ventas.xlsx"

with zipfile.ZipFile(base) as zin:
    names = zin.namelist()
    contents = {n: zin.read(n) for n in names}

rels_name = "ppt/charts/_rels/chart1.xml.rels"
rels = contents[rels_name].decode("utf-8")

# Localizar la relación al paquete embebido y convertirla en externa.
def to_external(m):
    tag = m.group(0)
    if "package" in tag or "oleObject" in tag or "embeddings" in tag.lower():
        tag = re.sub(r'Target="[^"]*"', lambda _: f'Target="{ext_target}"', tag)
        if "TargetMode" not in tag:
            tag = tag[:-2] + ' TargetMode="External"/>' if tag.endswith("/>") else tag
        else:
            tag = re.sub(r'TargetMode="[^"]*"', 'TargetMode="External"', tag)
    return tag

rels = re.sub(r"<Relationship\b[^>]*?/?>", to_external, rels)
contents[rels_name] = rels.encode("utf-8")
print("chart1.xml.rels reescrito:\n", rels)

with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:
    for n in names:
        zout.writestr(n, contents[n])

print("\nGenerado", out)
